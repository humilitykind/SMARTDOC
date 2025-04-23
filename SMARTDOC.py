# app.py

# ==================== IMPORTS ====================
from dotenv import load_dotenv
import os
import streamlit as st
import pandas as pd
import json
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings  # Updated import
from langchain_community.vectorstores import FAISS
from langchain.prompts import PromptTemplate
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain_community.chat_models import ChatOpenAI  # Updated import
from agentic_doc.parse import parse_documents
import tempfile
import mimetypes
from langchain.schema import Document as LC_Document
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
import warnings
from pptx import Presentation
from docx import Document 


# Suppress warnings
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=UserWarning)

# ==================== SETUP ====================
load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
VISION_AGENT_API_KEY = os.getenv("VISION_AGENT_API_KEY")

st.set_page_config(page_title="📚 Smart Multi-Doc AI Chatbot", page_icon="🧠")


# ==================== CUSTOM CSS ====================
st.markdown("""
<style>
    /* Main container */
    .stApp {
        background: linear-gradient(180deg, #00003b 0%, #000000 100%); /* Dark blue to black */
        color: #e2e8f0;
        font-family: 'Segoe UI', system-ui, sans-serif;
    }

    /* Sidebar */
    [data-testid="stSidebar"] {
        background: #1e293b !important;
        border-right: 1px solid #334155;
    }

    /* File uploader */
    [data-testid="stFileUploader"] {
        border: 2px dashed #475569 !important;
        border-radius: 10px !important;
        padding: 20px !important;
    }

    /* Buttons */
    .stButton>button {
        background: #3b82f6 !important;
        color: white !important;
        border-radius: 8px !important;
        padding: 8px 20px !important;
        transition: all 0.3s !important;
    }

    .stButton>button:hover {
        background: #2563eb !important;
        transform: translateY(-1px);
    }

    /* Chat messages */
    .stChatMessage {
        background: #1e293b !important;
        border-radius: 12px !important;
        padding: 16px !important;
        margin: 8px 0 !important;
        box-shadow: 0 2px 8px rgba(0,0,0,0.1);
    }

    /* Progress bar */
    .stProgress > div > div > div {
        background: linear-gradient(90deg, #3b82f6 0%, #6366f1 100%);
    }

    /* Expander */
    [data-testid="stExpander"] {
        background: #1e293b !important;
        border: 1px solid #334155 !important;
        border-radius: 8px !important;
    }

    /* Main heading */
    h1 {
        color: #ffffff !important; /* Very light blue */
    }
</style>
""", unsafe_allow_html=True)

# ==================== PROMPTS ====================
custom_template = """Given the following conversation and a follow up question, rephrase the follow up question to be a standalone question or multiple questions when needed, in its original language.
When analyzing JSON data:
- Look for UUID patterns (8-4-4-4-12 characters)
- Check both string and numeric ID fields
- Examine nested objects and arrays

Chat History:
{chat_history}
Follow Up Input: {question}
Standalone question:"""
CUSTOM_QUESTION_PROMPT = PromptTemplate.from_template(custom_template)

# ==================== UTILITIES ====================
def detect_file_type(file):
    """Enhanced file type detection"""
    mime_type, _ = mimetypes.guess_type(file.name)
    if file.name.lower().endswith(('.xlsx', '.xls')):
        return 'application/vnd.ms-excel'
    return mime_type or "application/octet-stream"

def process_single_file(file):
    """Robust file processor with enhanced error handling."""
    chunks = []
    df = None
    file_type = detect_file_type(file)
    _, ext = os.path.splitext(file.name)

    # Create temporary file with proper extension
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(file.read())
        tmp_path = tmp.name

    try:
        # PDF Processing (Accurate by default)
        if "pdf" in file_type:
            results = parse_documents([tmp_path])
            chunks = [chunk.text for chunk in results[0].chunks]
            return [LC_Document(
                page_content=c,
                metadata={
                    "source": file.name,
                    "file_type": "pdf",
                    "processing_mode": "accurate"
                }
            ) for c in chunks], None

        # Image Processing
        elif file_type.startswith('image/'):
            results = parse_documents([tmp_path])
            return [LC_Document(
                page_content=chunk.text,
                metadata={
                    "source": file.name,
                    "file_type": "image",
                    "processing_mode": "accurate"
                }
            ) for chunk in results[0].chunks], None

        # JSON Processing (Enhanced)
        elif file.name.lower().endswith('.json'):
            try:
                # print(f"Processing JSON file: {file.name}")  # Debugging: File name

                # Read the JSON file
                with open(tmp_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)

                # Debugging: Check the loaded JSON data
                # print("Loaded JSON Data:")
                # print(data)

                # Metadata for the JSON file
                metadata = {
                    "source": file.name,
                    "file_type": "json",
                    "entries": len(data) if isinstance(data, list) else 1
                }
                # print("Metadata:")
                # print(metadata)  # Debugging: Check metadata

                # Process JSON data based on its structure
                if isinstance(data, list):
                    # Debugging: Check the number of entries in the list
                    # print(f"JSON is a list with {len(data)} entries.")
                    chunks = [
                        f"JSON Entry {idx+1}:\n{json.dumps(item, indent=2)}"
                        for idx, item in enumerate(data)
                    ]
                elif isinstance(data, dict):
                    # Debugging: Check the dictionary structure
                    # print("JSON is a dictionary.")
                    content = json.dumps(data, indent=2)
                    # print("JSON Content:")
                    # print(content)  # Debugging: Check the JSON content
                    chunks = chunk_text(content)
                else:
                    raise ValueError("Unsupported JSON structure")

                # Debugging: Check the generated chunks
                # print("Generated Chunks:")
                # print(chunks)

                return [LC_Document(
                    page_content=c,
                    metadata=metadata
                ) for c in chunks], None

            except json.JSONDecodeError as e:
                st.error(f"Invalid JSON in {file.name}: {str(e)}")
                # print(f"JSON Decode Error: {str(e)}")  # Debugging: JSON decode error
                return [], None

            except Exception as e:
                st.error(f"Error processing JSON file: {str(e)}")
                # print(f"General Error: {str(e)}")  # Debugging: General error
                return [], None

        # Spreadsheet Processing (Excel/CSV)
        elif file.name.lower().endswith(('.xlsx', '.xls', '.csv')):
            try:
                # print(f"Processing file: {file.name}")  # Debugging: File name

                if file.name.endswith(('xlsx', 'xls')):
                    df = pd.read_excel(tmp_path)
                else:
                    df = pd.read_csv(tmp_path)

                # Debugging: Check the DataFrame
                # print("DataFrame Preview:")
                # print(df.head())

                content = df.to_markdown(index=False)

                # Debugging: Check the Markdown content
                # print("Markdown Content:")
                # print(content)

                chunks = chunk_text(content)

                # Debugging: Check the generated chunks
                # print("Generated Chunks:")
                # print(chunks)

                return [LC_Document(
                    page_content=c,
                    metadata={
                        "source": file.name,
                        "file_type": "spreadsheet",
                        "columns": str(df.columns.tolist())
                    }
                ) for c in chunks], df

            except Exception as e:
                st.error(f"Spreadsheet error: {str(e)}")
                # print(f"Error processing spreadsheet: {str(e)}")  # Debugging
                return [], None

        # Text File Processing
        elif file.name.lower().endswith('.txt'):
            try:
                with open(tmp_path, "r") as f:
                    text = f.read()
                chunks = chunk_text(text)
                return [LC_Document(
                    page_content=c,
                    metadata={
                        "source": file.name,
                        "file_type": "text"
                    }
                ) for c in chunks], None

            except Exception as e:
                st.error(f"Text file error: {str(e)}")
                return [], None

        # Word Document Processing (.docx)
        elif file.name.lower().endswith('.docx'):
            try:
                
                doc = Document(tmp_path)
                content = "\n".join([paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()])
                chunks = chunk_text(content)
                return [LC_Document(
                    page_content=c,
                    metadata={
                        "source": file.name,
                        "file_type": "docx"
                    }
                ) for c in chunks], None

            except Exception as e:
                st.error(f"Word document error: {str(e)}")
                return [], None

        # PowerPoint Presentation Processing (.pptx)
        elif file.name.lower().endswith('.pptx'):
            try:
                ppt = Presentation(tmp_path)
                content = "\n".join([
                    shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
                ])
                chunks = chunk_text(content)
                return [LC_Document(
                    page_content=c,
                    metadata={
                        "source": file.name,
                        "file_type": "pptx"
                    }
                ) for c in chunks], None

            except Exception as e:
                st.error(f"PowerPoint presentation error: {str(e)}")
                return [], None

        # Unsupported File Types
        else:
            st.warning(f"Unsupported file type: {file.name}")
            return [], None

    except Exception as e:
        st.error(f"Critical processing error: {str(e)}")
        return [], None

    finally:
        try:
            os.unlink(tmp_path)
        except Exception as e:
            st.error(f"Cleanup error: {str(e)}")

def chunk_text(text):
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", "。", " ", ""]
    )
    return splitter.split_text(text)

def build_vectorstore(documents):
    embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
    return FAISS.from_documents(documents, embeddings)

# ==================== PROCESSING PIPELINE ====================
def process_files(files):
    all_chunks = []
    dfs = {}
    
    with ThreadPoolExecutor(max_workers=4) as executor:
        futures = {executor.submit(process_single_file, file): file for file in files}
        
        for future in as_completed(futures):
            file = futures[future]
            chunks, df = future.result()
            if chunks:
                all_chunks.extend(chunks)
            if df is not None:
                dfs[file.name] = df

    return all_chunks, dfs

# ==================== CHAT ENGINE ====================
def build_conversation_chain(vectorstore):
    llm = ChatOpenAI(
        temperature=0.2,
        model_name="gpt-4o",
        max_tokens=2000
    )
    
    memory = ConversationBufferMemory(
        memory_key='chat_history',
        return_messages=True,
        output_key='answer'
    )

    return ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(search_kwargs={"k": 5}),
        condense_question_prompt=CUSTOM_QUESTION_PROMPT,
        memory=memory,
        max_tokens_limit=4000,
        return_source_documents=True
    )

# ==================== MAIN INTERFACE ====================
def main():
    # Hero Section
    col1, col2 = st.columns([3, 1])
    with col1:
        st.markdown("""
        <div style="padding: 20px 0 40px 0">
            <h1 style="font-size: 2.5rem; color: #3b82f6; margin-bottom: 0.5rem">
                🧠 SmartDoc Analyzer
            </h1>
            <p style="font-size: 1.1rem; color: #94a3b8">
                <em>Secure Next-gen document intelligence on any file type.</em>
            </p>
        </div>
        """, unsafe_allow_html=True)

    with st.sidebar:
        # File Uploader with custom styling
        uploaded_files = st.file_uploader(
            "📁 Upload Documents",
            type=["pdf", "csv", "xlsx", "txt", "png", "jpg", "jpeg", "json", "docx", "pptx"],
            accept_multiple_files=True,
            help="Supported Types: PDF, Images, Spreadsheets, Text, JSON, Word, PowerPoint"
        )

        # Progress Bar Setup
        progress_bar = st.progress(0)
        status_text = st.empty()

    # Processing Section
    if st.sidebar.button("🚀 Process Documents", use_container_width=True) and uploaded_files:
        with st.spinner("Hang tight, while we chomp on your documents ! ...") as first_spinner:
            total_files = len(uploaded_files)
            processed_files = 0
            all_chunks = []
            dfs = {}

            status_text.text(f"Processing 0/{total_files} files...")

            with ThreadPoolExecutor(max_workers=4) as executor:
                futures = {executor.submit(
                    process_single_file, 
                    file
                ): file for file in uploaded_files}

                for future in as_completed(futures):
                    processed_files += 1
                    progress = processed_files / total_files
                    progress_bar.progress(min(progress, 1.0))
                    status_text.text(f"Processing {processed_files}/{total_files} files...")

                    file = futures[future]
                    chunks, df = future.result()
                    if chunks:
                        all_chunks.extend(chunks)
                    if df is not None:
                        dfs[file.name] = df

            if all_chunks:
                first_spinner = None  # Stop the first spinner
                with st.spinner("Almost done, Swallowing now !..."):
                    vectorstore = build_vectorstore(all_chunks)
                    st.session_state.conversation = build_conversation_chain(vectorstore)
                    st.session_state.dataframes = dfs

                progress_bar.empty()
                status_text.success(f"✅ Processed {total_files} files successfully!")
            else:
                status_text.error("❌ Failed to extract meaningful content from files")

    # Chat Interface
    if "conversation" in st.session_state:
        st.markdown("### Document Chat 💬")

        # Display chat history
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []

        # User input
        question = st.chat_input("Ask about your documents:")

        if question:
            # Immediately show user question
            with st.chat_message("user"):
                st.markdown(question)

            # Generate response
            with st.spinner("Analyzing documents..."):
                try:
                    response = st.session_state.conversation.invoke({"question": question})
                    st.session_state.chat_history = response["chat_history"]

                    # Show assistant response
                    with st.chat_message("assistant"):
                        st.markdown(response["answer"])

                    # Source documents
                    with st.expander("🔍 View Source References"):
                        for doc in response.get("source_documents", []):
                            st.markdown(f"*Source:* {doc.metadata['source']}")
                            st.code(doc.page_content[:500] + "...", language="text")

                except Exception as e:
                    st.error(f"⚠ Error processing request: {str(e)}")

if __name__ == '__main__':
    main()